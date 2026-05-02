"""
generate_monthly_review.py — Monthly Review report for GLPG Risk Monitor.

PA contributions are alpha net of benchmark (REPORT_ALPHA_ATRIBUTION is alpha by design).
Benchmark-tracking LIVROs (_PA_BENCH_LIVROS) are separated from active totals.

Outputs in data/monthly-reviews/:
    {YYYY-MM}_monthly_review.html   — tabbed HTML, one tab per fund + market
    {YYYY-MM}_monthly_review.pptx   — one slide per fund (IDKAs combined), comment box per PM

Usage:
    python generate_monthly_review.py --month 2026-04
    python generate_monthly_review.py          # prompts for month
"""
from __future__ import annotations

import argparse
import html as html_lib
import json
import os
import sys
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Optional

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import pandas as pd

from glpg_fetch import read_sql
from month_bdays import month_bdays
from risk_config import (
    FUNDS, RAW_FUNDS, IDKA_FUNDS, ALL_FUNDS,
    FUND_ORDER, FUND_LABELS, _FUND_PA_KEY,
    _PA_BENCH_LIVROS, _FUND_PEERS_GROUP,
)

# ── Output directory ──────────────────────────────────────────────────────────
OUT_DIR = Path(__file__).parent / "data" / "monthly-reviews"
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ── Brand palette ─────────────────────────────────────────────────────────────
C_NAVY   = "#183B80"
C_BLUE   = "#2AADF5"
C_DARK   = "#133979"
C_WHITE  = "#FFFFFF"
C_GREY   = "#F5F7FA"
C_DGREY  = "#DEE2E9"
C_GREEN  = "#1E8C45"
C_RED    = "#C0392B"
C_LG     = "#D5F5E3"
C_LR     = "#FADBD8"

# ── Fund metadata ──────────────────────────────────────────────────────────────
_SHORT_TO_DESK: dict[str, str] = {cfg["short"]: name for name, cfg in ALL_FUNDS.items()}
_FUND_VAR_SOFT: dict[str, float] = {cfg["short"]: cfg["var_soft"] for cfg in ALL_FUNDS.values()}
_FUND_VAR_HARD: dict[str, float] = {cfg["short"]: cfg["var_hard"] for cfg in ALL_FUNDS.values()}

_BENCH_NAME: dict[str, str] = {
    "MACRO": "CDI",     "QUANT": "CDI",     "EVOLUTION": "CDI",
    "MACRO_Q": "CDI",   "ALBATROZ": "CDI",  "BALTRA": "CDI",
    "FRONTIER": "IBOV", "IDKA_3Y": "IDKA 3Y", "IDKA_10Y": "IDKA 10Y",
}
_IDKA_IDX_NAME: dict[str, str] = {
    "IDKA_3Y":  "IDKA_IPCA_3A",
    "IDKA_10Y": "IDKA_IPCA_10A",
}
_MONTH_PTBR: dict[int, str] = {
    1: "Janeiro", 2: "Fevereiro", 3: "Março",    4: "Abril",
    5: "Maio",    6: "Junho",     7: "Julho",     8: "Agosto",
    9: "Setembro",10: "Outubro",  11: "Novembro", 12: "Dezembro",
}

# LIVROs that are FX hedge collateral (mark-to-market noise, not alpha) —
# excluded from product-level top contributors/detractors. Matches daily
# pmovers_renderers._FX_HEDGE_LIVROS.
_FX_HEDGE_LIVROS: frozenset[str] = frozenset({
    "Caixa USD", "Cash USD", "Caixa USD Futures",
})
# Product names that are pure cash placeholders — excluded from product top-movers.
_NOISE_PRODUCTS: frozenset[str] = frozenset({"Cash USD"})

# Peers JSON — self-contained path inside Risk Monitor project.
# The file is maintained by the risk-data-collector / daily batch.
# Fallback: shared network path (may be unavailable when migrated).
_PEERS_JSON_LOCAL  = Path(__file__).parent / "data" / "peers_data.json"
_PEERS_JSON_REMOTE = Path(os.environ.get(
    "PEERS_DATA_PATH",
    r"\\fs02\FS_GALAPAGOS\Bloomberg\Quant\Claude_GLPG_Fetch\peers_data.json",
))

# HTML tab order (IDKAs merged into one tab)
_TAB_ORDER = ["MACRO", "QUANT", "EVOLUTION", "MACRO_Q", "ALBATROZ", "BALTRA", "FRONTIER", "IDKA"]
_TAB_LABELS = {
    "MACRO": "Macro", "QUANT": "Quantitativo", "EVOLUTION": "Evolution",
    "MACRO_Q": "Macro Q", "ALBATROZ": "Albatroz", "BALTRA": "Baltra",
    "FRONTIER": "Frontier", "IDKA": "IDKA",
}


# ── Data model ────────────────────────────────────────────────────────────────
@dataclass
class FundCtx:
    short: str
    label: str
    bench_name: str
    pa_by_livro: pd.DataFrame       # cols: livro, mtd_bps, ytd_bps, is_bench
    pa_products: pd.DataFrame       # cols: product, livro, classe, mtd_bps, ytd_bps
    mtd_alpha_bps: float = 0.0
    ytd_alpha_bps: float = 0.0
    fund_mtd_pct:  Optional[float] = None
    fund_ytd_pct:  Optional[float] = None
    bench_mtd_bps: float = 0.0
    bench_ytd_bps: float = 0.0
    var_avg_pct:   Optional[float] = None
    var_max_pct:   Optional[float] = None
    var_last_pct:  Optional[float] = None
    var_soft: float = 0.0
    var_hard: float = 0.0


# ── Date resolution ───────────────────────────────────────────────────────────
def _resolve_dates(ym: str) -> tuple[list[str], str, str, str, bool, str]:
    """Return (bdays, start, end, year_start, is_final, label_ptbr)."""
    y, m = int(ym[:4]), int(ym[5:7])
    bdays = month_bdays(ym)
    if not bdays:
        sys.exit(f"Nenhum dado encontrado para o mês {ym}.")
    start      = bdays[0]
    end        = bdays[-1]
    last_cal   = (pd.Timestamp(y, m, 1) + pd.offsets.MonthEnd(0)).date()
    is_final   = pd.Timestamp("today").date() > last_cal
    year_start = f"{y}-01-01"
    label      = f"{_MONTH_PTBR[m]} {y}"
    return bdays, start, end, year_start, is_final, label


# ── Data fetch ────────────────────────────────────────────────────────────────
def _fetch_pa(start: str, end: str, year_start: str) -> pd.DataFrame:
    """All-fund PA by (FUNDO, LIVRO, PRODUCT, CLASSE), MTD = start→end, YTD = year_start→end."""
    pa_keys = "','".join(_FUND_PA_KEY.values())
    df = read_sql(f"""
        SELECT "FUNDO", "LIVRO", "PRODUCT", "CLASSE",
               SUM(CASE WHEN "DATE" >= DATE '{start}'      AND "DATE" <= DATE '{end}'
                        THEN "DIA" ELSE 0 END) * 10000 AS mtd_bps,
               SUM(CASE WHEN "DATE" >= DATE '{year_start}' AND "DATE" <= DATE '{end}'
                        THEN "DIA" ELSE 0 END) * 10000 AS ytd_bps
        FROM q_models."REPORT_ALPHA_ATRIBUTION"
        WHERE "FUNDO" IN ('{pa_keys}')
          AND "DATE" >= DATE '{year_start}' - INTERVAL '5 days'
          AND "DATE" <= DATE '{end}'
        GROUP BY "FUNDO", "LIVRO", "PRODUCT", "CLASSE"
        HAVING ABS(SUM(CASE WHEN "DATE" >= DATE '{start}'      AND "DATE" <= DATE '{end}'      THEN "DIA" ELSE 0 END)) > 1e-8
            OR ABS(SUM(CASE WHEN "DATE" >= DATE '{year_start}' AND "DATE" <= DATE '{end}'      THEN "DIA" ELSE 0 END)) > 1e-8
    """)
    for c in ("mtd_bps", "ytd_bps"):
        df[c] = df[c].astype(float).fillna(0.0)
    return df


def _fetch_nav_series(year_start: str, end: str) -> pd.DataFrame:
    """Returns TRADING_DESK, VAL_DATE, SHARE (unit price for returns), NAV (total BRL for VaR %)."""
    desks = "','".join(ALL_FUNDS.keys())
    df = read_sql(f"""
        SELECT "TRADING_DESK", "VAL_DATE", "SHARE", "NAV"
        FROM "LOTE45"."LOTE_TRADING_DESKS_NAV_SHARE"
        WHERE "TRADING_DESK" IN ('{desks}')
          AND "VAL_DATE" >= DATE '{year_start}' - INTERVAL '15 days'
          AND "VAL_DATE" <= DATE '{end}'
        ORDER BY "TRADING_DESK", "VAL_DATE"
    """)
    df["VAL_DATE"] = pd.to_datetime(df["VAL_DATE"])
    df["SHARE"]    = df["SHARE"].astype(float)
    df["NAV"]      = df["NAV"].astype(float)
    return df


def _fetch_cdi(year_start: str, end: str, start: str) -> dict:
    df = read_sql(f"""
        SELECT
          SUM(CASE WHEN "DATE" >= DATE '{start}'      AND "DATE" <= DATE '{end}' THEN "VALUE" ELSE 0 END) * 10000 AS mtd_bps,
          SUM(CASE WHEN "DATE" >= DATE '{year_start}' AND "DATE" <= DATE '{end}' THEN "VALUE" ELSE 0 END) * 10000 AS ytd_bps
        FROM public."ECO_INDEX"
        WHERE "INSTRUMENT" = 'CDI' AND "FIELD" = 'YIELD'
          AND "DATE" >= DATE '{year_start}' AND "DATE" <= DATE '{end}'
    """)
    if df.empty:
        return {"mtd_bps": 0.0, "ytd_bps": 0.0}
    return {"mtd_bps": float(df["mtd_bps"].iloc[0] or 0),
            "ytd_bps": float(df["ytd_bps"].iloc[0] or 0)}


def _fetch_ibov(year_start: str, end: str, start: str) -> dict:
    df = read_sql(f"""
        SELECT "DATE", "CLOSE"
        FROM public."EQUITIES_PRICES"
        WHERE "INSTRUMENT" = 'IBOV'
          AND "DATE" >= DATE '{year_start}' - INTERVAL '15 days'
          AND "DATE" <= DATE '{end}'
        ORDER BY "DATE"
    """)
    if df.empty:
        return {"mtd_bps": 0.0, "ytd_bps": 0.0}
    df["DATE"] = pd.to_datetime(df["DATE"])
    df = df.drop_duplicates("DATE").set_index("DATE").sort_index()

    def _ret(from_before: str, to_date: str) -> float:
        prev = df[df.index < pd.Timestamp(from_before)]
        curr = df[df.index <= pd.Timestamp(to_date)]
        if prev.empty or curr.empty:
            return 0.0
        p0 = float(prev["CLOSE"].iloc[-1])
        p1 = float(curr["CLOSE"].iloc[-1])
        return (p1 / p0 - 1) * 10000 if p0 else 0.0

    return {
        "mtd_bps": _ret(start, end),
        "ytd_bps": _ret(year_start, end),
    }


def _fetch_idka_returns(year_start: str, end: str, start: str) -> dict[str, dict]:
    idx_names = "','".join(_IDKA_IDX_NAME.values())
    df = read_sql(f"""
        SELECT "INSTRUMENT", "DATE", "VALUE"
        FROM public."ECO_INDEX"
        WHERE "INSTRUMENT" IN ('{idx_names}') AND "FIELD" = 'INDEX'
          AND "DATE" >= DATE '{year_start}' - INTERVAL '15 days'
          AND "DATE" <= DATE '{end}'
        ORDER BY "INSTRUMENT", "DATE"
    """)
    result: dict[str, dict] = {}
    for short, idx_name in _IDKA_IDX_NAME.items():
        if df.empty:
            result[short] = {"mtd_bps": 0.0, "ytd_bps": 0.0}
            continue
        sub = df[df["INSTRUMENT"] == idx_name].copy()
        sub["DATE"] = pd.to_datetime(sub["DATE"])
        sub = sub.drop_duplicates("DATE").set_index("DATE").sort_index()

        def _ret(from_before: str, to_date: str, s=sub) -> float:
            prev = s[s.index < pd.Timestamp(from_before)]
            curr = s[s.index <= pd.Timestamp(to_date)]
            if prev.empty or curr.empty:
                return 0.0
            p0 = float(prev["VALUE"].iloc[-1])
            p1 = float(curr["VALUE"].iloc[-1])
            return (p1 / p0 - 1) * 10000 if p0 else 0.0

        result[short] = {
            "mtd_bps": _ret(start, end),
            "ytd_bps": _ret(year_start, end),
        }
    return result


def _nav_return(nav_df: pd.DataFrame, short: str,
                start: str, end: str, year_start: str) -> dict:
    desk = _SHORT_TO_DESK.get(short)
    if not desk:
        return {"mtd_pct": None, "ytd_pct": None}
    sub = nav_df[nav_df["TRADING_DESK"] == desk].set_index("VAL_DATE").sort_index()["SHARE"]
    if sub.empty:
        return {"mtd_pct": None, "ytd_pct": None}

    def _ret(from_before: str, to_date: str) -> Optional[float]:
        prev = sub[sub.index < pd.Timestamp(from_before)]
        curr = sub[sub.index <= pd.Timestamp(to_date)]
        if prev.empty or curr.empty:
            return None
        base = float(prev.iloc[-1])
        return (float(curr.iloc[-1]) / base - 1) * 100 if base else None

    return {
        "mtd_pct": _ret(start, end),
        "ytd_pct": _ret(year_start, end),
    }


def _fetch_var_period(start: str, end: str, nav_df: pd.DataFrame) -> dict[str, dict]:
    out: dict[str, dict] = {}

    # RPM funds
    if FUNDS:
        level_clause = " OR ".join(
            f"(\"TRADING_DESK\" = '{td}' AND \"LEVEL\" = {cfg['level']})"
            for td, cfg in FUNDS.items()
        )
        try:
            df_rpm = read_sql(f"""
                SELECT "TRADING_DESK", "VAL_DATE", SUM("PARAMETRIC_VAR") AS var_abs
                FROM "LOTE45"."LOTE_FUND_STRESS_RPM"
                WHERE "VAL_DATE" >= DATE '{start}' AND "VAL_DATE" <= DATE '{end}'
                  AND ({level_clause})
                GROUP BY "TRADING_DESK", "VAL_DATE"
                ORDER BY "TRADING_DESK", "VAL_DATE"
            """)
            df_rpm["VAL_DATE"] = pd.to_datetime(df_rpm["VAL_DATE"])
            for td, cfg in FUNDS.items():
                short = cfg["short"]
                sub = df_rpm[df_rpm["TRADING_DESK"] == td].set_index("VAL_DATE").sort_index()
                if sub.empty:
                    continue
                nav_sub = nav_df[nav_df["TRADING_DESK"] == td].set_index("VAL_DATE").sort_index()["NAV"]
                nav_aligned = nav_sub.reindex(sub.index.union(nav_sub.index), method="ffill").reindex(sub.index)
                sub["var_pct"] = sub["var_abs"].abs() / nav_aligned * 100
                sub = sub.dropna(subset=["var_pct"])
                if not sub.empty:
                    out[short] = {"avg": float(sub["var_pct"].mean()),
                                  "max": float(sub["var_pct"].max()),
                                  "last": float(sub["var_pct"].iloc[-1])}
        except Exception:
            pass

    # RAW funds
    if RAW_FUNDS:
        tds = "','".join(RAW_FUNDS.keys())
        try:
            # TREE='Main' filter: BALTRA has 3 TREEs (Main, Main_Macro_Gestores,
            # Main_Macro_Ativos), each containing the full fund — without filter
            # the SUM triplicates. Other RAW_FUNDS only have Main → no-op.
            df_raw = read_sql(f"""
                SELECT "TRADING_DESK", "VAL_DATE", SUM("PVAR1DAY") AS var_abs
                FROM "LOTE45"."LOTE_FUND_STRESS"
                WHERE "VAL_DATE" >= DATE '{start}' AND "VAL_DATE" <= DATE '{end}'
                  AND "TRADING_DESK" IN ('{tds}')
                  AND "TREE" = 'Main'
                GROUP BY "TRADING_DESK", "VAL_DATE"
                ORDER BY "TRADING_DESK", "VAL_DATE"
            """)
            df_raw["VAL_DATE"] = pd.to_datetime(df_raw["VAL_DATE"])
            for td, cfg in RAW_FUNDS.items():
                short = cfg["short"]
                sub = df_raw[df_raw["TRADING_DESK"] == td].set_index("VAL_DATE").sort_index()
                if sub.empty:
                    continue
                nav_sub = nav_df[nav_df["TRADING_DESK"] == td].set_index("VAL_DATE").sort_index()["NAV"]
                nav_aligned = nav_sub.reindex(sub.index.union(nav_sub.index), method="ffill").reindex(sub.index)
                sub["var_pct"] = sub["var_abs"].abs() / nav_aligned * 100
                sub = sub.dropna(subset=["var_pct"])
                if not sub.empty:
                    out[short] = {"avg": float(sub["var_pct"].mean()),
                                  "max": float(sub["var_pct"].max()),
                                  "last": float(sub["var_pct"].iloc[-1])}
        except Exception:
            pass

    # IDKA funds (RELATIVE_VAR_PCT already fractional BVaR)
    if IDKA_FUNDS:
        tds = "','".join(IDKA_FUNDS.keys())
        try:
            df_idka = read_sql(f"""
                SELECT "TRADING_DESK", "VAL_DATE", SUM("RELATIVE_VAR_PCT") AS bvar_raw
                FROM "LOTE45"."LOTE_PARAMETRIC_VAR_TABLE"
                WHERE "VAL_DATE" >= DATE '{start}' AND "VAL_DATE" <= DATE '{end}'
                  AND "TRADING_DESK" IN ('{tds}')
                  AND "BOOKS"::text = '{{*}}'
                GROUP BY "TRADING_DESK", "VAL_DATE"
                ORDER BY "TRADING_DESK", "VAL_DATE"
            """)
            df_idka["VAL_DATE"] = pd.to_datetime(df_idka["VAL_DATE"])
            for td, cfg in IDKA_FUNDS.items():
                short = cfg["short"]
                sub = df_idka[df_idka["TRADING_DESK"] == td].set_index("VAL_DATE").sort_index()
                if sub.empty:
                    continue
                # abs() to mirror the RPM/RAW path convention (line 311, 344). The prior
                # `-bvar_raw` relied on RELATIVE_VAR_PCT always being negative-signed —
                # any positive entry would silently flip max() to under-report risk.
                sub["var_pct"] = sub["bvar_raw"].abs() * 100
                out[short] = {"avg": float(sub["var_pct"].mean()),
                              "max": float(sub["var_pct"].max()),
                              "last": float(sub["var_pct"].iloc[-1])}
        except Exception:
            pass

    return out


def _fetch_market_moves(start: str, end: str, year_start: str) -> list[dict]:
    rows = []

    def _eco_index_ret(instrument: str, field: str, label: str, as_pct: bool = True) -> Optional[dict]:
        try:
            df = read_sql(f"""
                SELECT "DATE", "VALUE"
                FROM public."ECO_INDEX"
                WHERE "INSTRUMENT" = '{instrument}' AND "FIELD" = '{field}'
                  AND "DATE" >= DATE '{year_start}' - INTERVAL '15 days'
                  AND "DATE" <= DATE '{end}'
                ORDER BY "DATE"
            """)
            if df.empty:
                return None
            df["DATE"] = pd.to_datetime(df["DATE"])
            df = df.set_index("DATE").sort_index()
            prev_mtd = df[df.index < pd.Timestamp(start)]
            prev_ytd = df[df.index < pd.Timestamp(year_start)]
            curr     = df[df.index <= pd.Timestamp(end)]
            if prev_mtd.empty or curr.empty:
                return None
            s_mtd = float(prev_mtd["VALUE"].iloc[-1])
            s_ytd = float(prev_ytd["VALUE"].iloc[-1]) if not prev_ytd.empty else s_mtd
            e     = float(curr["VALUE"].iloc[-1])
            if as_pct:
                return {"label": label, "val_start": s_mtd, "val_end": e,
                        "chg_mtd": (e/s_mtd - 1)*100 if s_mtd else 0,
                        "chg_ytd": (e/s_ytd - 1)*100 if s_ytd else 0, "unit": "%"}
            else:
                return {"label": label, "val_start": s_mtd, "val_end": e,
                        "chg_mtd": (e - s_mtd)*10000, "chg_ytd": (e - s_ytd)*10000, "unit": "bps"}
        except Exception:
            return None

    def _price_ret(instrument: str, field: str, label: str) -> Optional[dict]:
        try:
            df = read_sql(f"""
                SELECT "DATE", "{field}" AS val
                FROM public."EQUITIES_PRICES"
                WHERE "INSTRUMENT" = '{instrument}'
                  AND "DATE" >= DATE '{year_start}' - INTERVAL '15 days'
                  AND "DATE" <= DATE '{end}'
                ORDER BY "DATE"
            """)
            if df.empty:
                return None
            df["DATE"] = pd.to_datetime(df["DATE"])
            df = df.drop_duplicates("DATE").set_index("DATE").sort_index()
            prev_mtd = df[df.index < pd.Timestamp(start)]
            prev_ytd = df[df.index < pd.Timestamp(year_start)]
            curr     = df[df.index <= pd.Timestamp(end)]
            if prev_mtd.empty or curr.empty:
                return None
            s_mtd = float(prev_mtd["val"].iloc[-1])
            s_ytd = float(prev_ytd["val"].iloc[-1]) if not prev_ytd.empty else s_mtd
            e     = float(curr["val"].iloc[-1])
            return {"label": label, "val_start": s_mtd, "val_end": e,
                    "chg_mtd": (e/s_mtd - 1)*100 if s_mtd else 0,
                    "chg_ytd": (e/s_ytd - 1)*100 if s_ytd else 0, "unit": "%"}
        except Exception:
            return None

    candidates = [
        ("IBOV",           None,      "IBOV"),
        ("CDI",            "INDEX",   "CDI (acum)"),
        ("IDKA_IPCA_3A",   "INDEX",   "IDKA 3Y"),
        ("IDKA_IPCA_10A",  "INDEX",   "IDKA 10Y"),
    ]
    for inst, fld, lbl in candidates:
        if fld is None:
            r = _price_ret(inst, "CLOSE", lbl)
        else:
            r = _eco_index_ret(inst, fld, lbl)
        if r:
            rows.append(r)

    return rows


# ── Data assembly ─────────────────────────────────────────────────────────────
def _build_fund_ctx(
    short: str,
    pa_df: pd.DataFrame,
    nav_df: pd.DataFrame,
    benchmarks: dict,
    var_stats: dict[str, dict],
    start: str, end: str, year_start: str,
) -> FundCtx:
    label      = FUND_LABELS.get(short, short)
    bench_name = _BENCH_NAME.get(short, "CDI")
    pa_key     = _FUND_PA_KEY.get(short, short)
    bench_livros = _PA_BENCH_LIVROS.get(pa_key, set())

    sub = pa_df[pa_df["FUNDO"] == pa_key].copy()

    # --- PA by CLASSE (category) — exclude bench livros, merge Caixa/Custos classes ---
    _NOISE_CLASSES = frozenset({"Caixa", "Custos"})
    sub_active = sub[~sub["LIVRO"].isin(bench_livros)]

    by_cat = (
        sub_active
        .groupby("CLASSE")[["mtd_bps", "ytd_bps"]]
        .sum()
        .reset_index()
        .rename(columns={"CLASSE": "livro"})
    )
    # Merge Caixa + Custos classes into one row
    noise_mask = by_cat["livro"].isin(_NOISE_CLASSES)
    if noise_mask.any():
        merged = pd.DataFrame([{
            "livro": "Caixa & Custos",
            "mtd_bps": float(by_cat.loc[noise_mask, "mtd_bps"].sum()),
            "ytd_bps": float(by_cat.loc[noise_mask, "ytd_bps"].sum()),
            "is_bench": True,
        }])
        by_cat = pd.concat([by_cat[~noise_mask], merged], ignore_index=True)
    by_cat["is_bench"] = by_cat["livro"] == "Caixa & Custos"
    by_cat = by_cat.sort_values("mtd_bps", ascending=False).reset_index(drop=True)

    by_livro = by_cat  # stored as pa_by_livro for display
    active   = by_cat[~by_cat["is_bench"]]
    mtd_alpha = float(active["mtd_bps"].sum())
    ytd_alpha = float(active["ytd_bps"].sum())

    # --- PA by product (top contributors/detractors; exclude noise classes
    #     + FX hedge collateral livros + Cash USD product). Matches daily
    #     pmovers filtering. ---
    sub_products = sub_active[
        (~sub_active["CLASSE"].isin(_NOISE_CLASSES))
        & (~sub_active["LIVRO"].isin(_FX_HEDGE_LIVROS))
        & (~sub_active["PRODUCT"].isin(_NOISE_PRODUCTS))
    ]
    products = (
        sub_products
           .groupby(["PRODUCT", "LIVRO", "CLASSE"])[["mtd_bps", "ytd_bps"]]
           .sum()
           .reset_index()
           .rename(columns={"PRODUCT": "product", "LIVRO": "livro", "CLASSE": "classe"})
    )

    # --- Fund NAV return ---
    nav_ret = _nav_return(nav_df, short, start, end, year_start)

    # --- Benchmark returns ---
    if bench_name == "CDI":
        bench_ret = benchmarks.get("cdi", {"mtd_bps": 0.0, "ytd_bps": 0.0})
    elif bench_name == "IBOV":
        bench_ret = benchmarks.get("ibov", {"mtd_bps": 0.0, "ytd_bps": 0.0})
    elif short == "IDKA_3Y":
        bench_ret = benchmarks.get("idka", {}).get("IDKA_3Y", {"mtd_bps": 0.0, "ytd_bps": 0.0})
    elif short == "IDKA_10Y":
        bench_ret = benchmarks.get("idka", {}).get("IDKA_10Y", {"mtd_bps": 0.0, "ytd_bps": 0.0})
    else:
        bench_ret = {"mtd_bps": 0.0, "ytd_bps": 0.0}

    # --- VaR stats ---
    vs = var_stats.get(short, {})

    return FundCtx(
        short=short, label=label, bench_name=bench_name,
        pa_by_livro=by_livro, pa_products=products,
        mtd_alpha_bps=mtd_alpha, ytd_alpha_bps=ytd_alpha,
        fund_mtd_pct=nav_ret["mtd_pct"], fund_ytd_pct=nav_ret["ytd_pct"],
        bench_mtd_bps=float(bench_ret.get("mtd_bps", 0)),
        bench_ytd_bps=float(bench_ret.get("ytd_bps", 0)),
        var_avg_pct=vs.get("avg"),   var_max_pct=vs.get("max"),
        var_last_pct=vs.get("last"),
        var_soft=_FUND_VAR_SOFT.get(short, 0.0),
        var_hard=_FUND_VAR_HARD.get(short, 0.0),
    )


# ── Formatting helpers ────────────────────────────────────────────────────────
def _fbps(v: float, decimals: int = 1) -> str:
    """Format bps with sign and BR decimal. e.g. +45.3 bps"""
    sign = "+" if v >= 0 else ""
    return f"{sign}{v:.{decimals}f} bps".replace(".", ",")


def _fpct(v: float, decimals: int = 2) -> str:
    """Format % with sign and BR decimal. e.g. +1,23%"""
    sign = "+" if v >= 0 else ""
    return f"{sign}{v:.{decimals}f}%".replace(".", ",")


def _css_sign(v: float) -> str:
    return "pos" if v >= 0 else "neg"


# ── Peers data ────────────────────────────────────────────────────────────────
def _load_peers_data() -> dict:
    """Load peers snapshot. Tries local copy first, then remote share."""
    for path in [_PEERS_JSON_LOCAL, _PEERS_JSON_REMOTE]:
        try:
            if path.exists():
                data = json.loads(path.read_text(encoding="utf-8"))
                return data if "val_date" in data else data.get("latest", {})
        except Exception:
            continue
    return {}


def _get_peer_list(short: str, peers_data: dict) -> list[dict]:
    """Peer fund dicts for a fund. Includes our fund (is_fund=True)."""
    group_key = _FUND_PEERS_GROUP.get(short)
    if not group_key or not peers_data:
        return []
    g = peers_data.get("groups", {}).get(group_key, {})
    return list(g.get("peers", []))


def _draw_peer_strip(ax, our_pct: float, tick_labels: list[str],
                     left_label: str, right_label: str) -> None:
    """Draw a single peer strip (matches daily report's mkStrip SVG):
    thin blue axis line with 5 ticks, diamond at our_pct (0–100) coloured by
    temperature (red worst → green best). Axis x ranges 0–100."""
    BLUE   = "#5aa3e8"
    BLUE_M = (90/255, 163/255, 232/255, 0.35)
    pc = max(0.0, min(100.0, our_pct))

    if pc < 20:   fill, edge = "#C0392B", "#7a1d12"
    elif pc < 40: fill, edge = "#E67E22", "#8a4a12"
    elif pc < 60: fill, edge = "#F4D03F", "#8a7a12"
    elif pc < 80: fill, edge = "#58D68D", "#1f7a3a"
    else:         fill, edge = "#1E8C45", "#0d4a22"

    # Halo + main blue line
    ax.plot([0, 100], [0.5, 0.5], color=BLUE_M, linewidth=6,
            solid_capstyle="round", zorder=2)
    ax.plot([0, 100], [0.5, 0.5], color=BLUE, linewidth=1.6,
            solid_capstyle="round", zorder=3)

    # End caps
    ax.scatter([0, 100], [0.5, 0.5], s=20, color=BLUE, zorder=4)

    # Tick marks at 0/25/50/75/100
    for i, t in enumerate([0, 25, 50, 75, 100]):
        h = 0.18 if i in (0, 4) else 0.22
        ax.plot([t, t], [0.5 - h, 0.5 + h], color=BLUE,
                linewidth=1.0, alpha=0.75, zorder=4)
        if tick_labels and i < len(tick_labels):
            ax.text(t, 0.10, tick_labels[i], ha="center", va="top",
                    fontsize=7, color="#888", family="monospace")

    # Diamond (temperature-coloured)
    ax.scatter([pc], [0.5], marker="D", s=110, color=fill,
               edgecolors=edge, linewidth=1.3, zorder=6)
    # Inner glassy highlight
    ax.scatter([pc], [0.5], marker="D", s=18,
               color=(1, 1, 1, 0.55), zorder=7)

    # Left label (outside axis)
    ax.text(-3, 0.5, left_label, ha="right", va="center",
            fontsize=8.5, fontweight="600", color="#222")
    # Right label
    ax.text(103, 0.5, right_label, ha="left", va="center",
            fontsize=8.5, fontweight="600", color="#222")

    ax.set_xlim(-32, 132)
    ax.set_ylim(-0.05, 1.05)
    ax.set_xticks([])
    ax.set_yticks([])
    for sp in ("top", "right", "left", "bottom"):
        ax.spines[sp].set_visible(False)
    ax.set_facecolor("white")


def _peers_chart(peers: list[dict], fund_label: str) -> BytesIO:
    """All-in-one peers figure: ranked bars + 2 strips per window
    (ordinal P0–P100 + retorno) matching the daily report pattern, + 2 scatters."""
    if not peers:
        return BytesIO()

    def _sorted_window(window: str) -> tuple[list[str], list[float], list[bool]]:
        rows = sorted(peers, key=lambda p: -(p.get(window) or 0))
        names  = [p["name"][:32] for p in rows]
        vals   = [(p.get(window) or 0) * 100 for p in rows]
        is_us  = [bool(p.get("is_fund")) for p in rows]
        return names, vals, is_us

    n = len(peers)
    # 4 rows: bar / ordinal strip / return strip / scatter
    height_ratios = [max(n * 0.28 + 0.6, 2.5), 0.55, 0.55, 2.2]
    fig = plt.figure(figsize=(13, sum(height_ratios) + 0.5),
                     facecolor="white")
    gs = fig.add_gridspec(4, 2, height_ratios=height_ratios,
                          hspace=0.75, wspace=0.30,
                          left=0.12, right=0.97, top=0.97, bottom=0.03)

    for col, window in enumerate(["MTD", "YTD"]):
        names, vals, is_us = _sorted_window(window)
        our_idx = next((i for i, f in enumerate(is_us) if f), 0)
        our_pct = (n - 1 - our_idx) / max(n - 1, 1) * 100

        # ── Row 0: sorted horizontal bar chart ────────────────────────────────
        ax_bar = fig.add_subplot(gs[0, col])
        bar_colors = [C_NAVY if u else "#A8C4E8" for u in is_us]
        ax_bar.barh(range(n), vals, color=bar_colors, height=0.62, zorder=3)
        ax_bar.set_yticks(range(n))
        ax_bar.set_yticklabels(
            [f"★ {nm}" if u else nm for nm, u in zip(names, is_us)],
            fontsize=7,
        )
        ax_bar.invert_yaxis()
        ax_bar.axvline(0, color="#bbb", linewidth=0.8)
        ax_bar.set_xlabel(f"{window} retorno %", fontsize=8)
        ax_bar.set_title(f"Ranking {window}", fontsize=9.5, fontweight="bold",
                         color=C_DARK)
        ax_bar.spines["top"].set_visible(False)
        ax_bar.spines["right"].set_visible(False)
        ax_bar.tick_params(axis="y", length=0)
        ax_bar.set_facecolor("white")
        max_abs = max(abs(v) for v in vals) if vals else 1
        for i, (val, u) in enumerate(zip(vals, is_us)):
            x = val + max_abs * 0.03 if val >= 0 else val - max_abs * 0.03
            ha = "left" if val >= 0 else "right"
            ax_bar.text(x, i, f"{val:+.2f}%", va="center", ha=ha,
                        fontsize=7.5 if u else 6.5,
                        fontweight="bold" if u else "normal",
                        color=C_NAVY if u else "#444")

        # Interpolated percentile reference lines: worst, P25, P50, P75, best
        sorted_vals = sorted(vals)
        q_refs = [
            (sorted_vals[0],            "#C0392B", "pior",  "--"),
            (sorted_vals[max(0, n//4)], "#E67E22", "P25",   ":"),
            (sorted_vals[max(0, n//2)], "#888888", "P50",   "-"),
            (sorted_vals[min(n-1, 3*n//4)], "#58D68D", "P75", ":"),
            (sorted_vals[-1],           "#1E8C45", "top",   "--"),
        ]
        for qval, qcol, qlbl, qls in q_refs:
            ax_bar.axvline(qval, color=qcol, linewidth=0.8, linestyle=qls,
                           alpha=0.7, zorder=2)
            ax_bar.text(qval, -0.7, qlbl, ha="center", va="top",
                        fontsize=6, color=qcol, alpha=0.9)

        # ── Rows 1–2: 2 strips per window (ordinal + retorno) ─────────────────
        # Mirrors daily report's per-fund peers chart: thin blue axis line,
        # 5 ticks (P0/P25/P50/P75/P100), diamond marker at our position
        # coloured by temperature (red worst → green best).
        ord_pct = (n - 1 - our_idx) / max(n - 1, 1) * 100
        sorted_asc = sorted(vals)
        p0, p100 = sorted_asc[0], sorted_asc[-1]

        def _q(p: float) -> float:
            if n <= 1:
                return sorted_asc[0]
            pos = p * (n - 1)
            lo, hi = int(np.floor(pos)), int(np.ceil(pos))
            return sorted_asc[lo] + (sorted_asc[hi] - sorted_asc[lo]) * (pos - lo)

        p25, p50, p75 = _q(0.25), _q(0.50), _q(0.75)
        our_val = vals[our_idx]
        val_pct = (our_val - p0) / (p100 - p0) * 100 if p100 > p0 else 50

        ax_ord = fig.add_subplot(gs[1, col])
        _draw_peer_strip(
            ax_ord, our_pct=ord_pct,
            tick_labels=["P0", "P25", "P50", "P75", "P100"],
            left_label=f"{window} · ordinal",
            right_label=f"P{ord_pct:.0f} · #{our_idx + 1}/{n}",
        )

        ax_ret = fig.add_subplot(gs[2, col])
        _draw_peer_strip(
            ax_ret, our_pct=val_pct,
            tick_labels=[f"{v:+.1f}%" for v in (p0, p25, p50, p75, p100)],
            left_label=f"{window} · retorno",
            right_label=f"{our_val:+.2f}%",
        )

    # ── Row 3 col 0: MTD vs YTD scatter ───────────────────────────────────────
    ax_sc1 = fig.add_subplot(gs[3, 0])
    for p in peers:
        xv = (p.get("MTD") or 0) * 100
        yv = (p.get("YTD") or 0) * 100
        if p.get("is_fund"):
            ax_sc1.scatter(xv, yv, color=C_NAVY, s=90, marker="*", zorder=5)
            ax_sc1.annotate(fund_label, (xv, yv), xytext=(5, 5),
                            textcoords="offset points", fontsize=7, color=C_NAVY)
        else:
            ax_sc1.scatter(xv, yv, color="#A8C4E8", s=22, alpha=0.8, zorder=3)
    ax_sc1.axhline(0, color="#ddd", lw=0.7)
    ax_sc1.axvline(0, color="#ddd", lw=0.7)
    ax_sc1.set_xlabel("MTD retorno %", fontsize=8)
    ax_sc1.set_ylabel("YTD retorno %", fontsize=8)
    ax_sc1.set_title("MTD vs YTD (consistência)", fontsize=9, fontweight="bold")
    ax_sc1.spines["top"].set_visible(False)
    ax_sc1.spines["right"].set_visible(False)
    ax_sc1.set_facecolor("white")

    # ── Row 3 col 1: Vol vs YTD scatter ───────────────────────────────────────
    ax_sc2 = fig.add_subplot(gs[3, 1])
    for p in peers:
        xv = (p.get("Vol") or 0) * 100
        yv = (p.get("YTD") or 0) * 100
        if p.get("is_fund"):
            ax_sc2.scatter(xv, yv, color=C_NAVY, s=90, marker="*", zorder=5)
            ax_sc2.annotate(fund_label, (xv, yv), xytext=(5, 5),
                            textcoords="offset points", fontsize=7, color=C_NAVY)
        else:
            ax_sc2.scatter(xv, yv, color="#A8C4E8", s=22, alpha=0.8, zorder=3)
    ax_sc2.axhline(0, color="#ddd", lw=0.7)
    ax_sc2.set_xlabel("Volatilidade % (12M)", fontsize=8)
    ax_sc2.set_ylabel("YTD retorno %", fontsize=8)
    ax_sc2.set_title("Risco × Retorno (YTD)", fontsize=9, fontweight="bold")
    ax_sc2.spines["top"].set_visible(False)
    ax_sc2.spines["right"].set_visible(False)
    ax_sc2.set_facecolor("white")

    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=130, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _html_peers_section(peers: list[dict], fund_label: str = "") -> str:
    if not peers:
        return ""
    # Ranking table sorted by MTD
    rows = sorted(peers, key=lambda p: -(p.get("MTD") or 0))
    trs = ""
    for rank, p in enumerate(rows, 1):
        is_us = p.get("is_fund", False)
        style = 'style="font-weight:700;background:#EEF2FF"' if is_us else ""
        star  = "★ " if is_us else ""
        trs += (
            f'<tr {style}><td>{rank}</td>'
            f'<td>{html_lib.escape(star + p["name"])}</td>'
            f'<td class="{_css_sign(p.get("MTD",0))}">{(p.get("MTD") or 0)*100:+.2f}%</td>'
            f'<td class="{_css_sign(p.get("YTD",0))}">{(p.get("YTD") or 0)*100:+.2f}%</td>'
            f'<td class="{_css_sign(p.get("12M",0))}">{(p.get("12M") or 0)*100:+.2f}%</td>'
            f'<td class="neutral">{(p.get("Vol") or 0)*100:.2f}%</td>'
            f'</tr>'
        )
    table_html = (
        '<table class="rt" style="max-width:800px">'
        '<thead><tr><th>#</th><th>Fundo</th>'
        '<th>MTD</th><th>YTD</th><th>12M</th><th>Vol</th></tr></thead>'
        f'<tbody>{trs}</tbody></table>'
    )
    # Embed the peers chart (bar + percentile strip + scatter) as base64 PNG
    chart_html = ""
    try:
        buf = _peers_chart(peers, fund_label or "★")
        if buf.getvalue():
            import base64
            b64 = base64.b64encode(buf.getvalue()).decode()
            chart_html = (
                f'<div style="margin-top:14px">'
                f'<img src="data:image/png;base64,{b64}" '
                f'style="max-width:100%;border-radius:6px;box-shadow:0 1px 4px rgba(0,0,0,.1)">'
                f'</div>'
            )
    except Exception:
        pass
    return (
        '<div class="section-title">Peers</div>'
        + table_html
        + chart_html
    )


# ── HTML CSS ──────────────────────────────────────────────────────────────────
_CSS = f"""
  :root {{
    --navy: {C_NAVY}; --blue: {C_BLUE}; --dark: {C_DARK};
    --green: {C_GREEN}; --red: {C_RED};
    --grey: {C_GREY}; --dgrey: {C_DGREY};
  }}
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: 'Segoe UI', Arial, sans-serif; background: #f0f2f5;
          color: #222; font-size: 13px; }}
  .header {{
    background: var(--navy); color: white; padding: 14px 24px;
    display: flex; align-items: center; justify-content: space-between;
    position: sticky; top: 0; z-index: 100; box-shadow: 0 2px 8px rgba(0,0,0,.25);
  }}
  .header-title {{ font-size: 18px; font-weight: 700; letter-spacing: .5px; }}
  .header-sub   {{ font-size: 12px; opacity: .85; }}
  .header-badge {{
    background: {C_BLUE}; color: white; padding: 4px 12px;
    border-radius: 12px; font-size: 11px; font-weight: 600;
  }}

  /* Tabs */
  .tab-bar {{
    background: var(--dark); display: flex; overflow-x: auto;
    padding: 0 16px; gap: 2px; position: sticky; top: 53px; z-index: 99;
  }}
  .tab-btn {{
    background: transparent; border: none; color: rgba(255,255,255,.65);
    padding: 10px 18px; cursor: pointer; font-size: 13px; font-weight: 500;
    white-space: nowrap; border-bottom: 3px solid transparent;
    transition: color .15s, border-color .15s;
  }}
  .tab-btn:hover  {{ color: white; border-color: rgba(255,255,255,.3); }}
  .tab-btn.active {{ color: white; border-color: {C_BLUE}; }}
  .tab-content {{ display: none; padding: 18px 20px; }}
  .tab-content.active {{ display: block; }}

  /* Cards */
  .metric-row {{ display: flex; gap: 12px; margin-bottom: 14px; flex-wrap: wrap; }}
  .metric-card {{
    background: white; border-radius: 8px; padding: 12px 18px;
    flex: 1; min-width: 160px; border-top: 3px solid var(--navy);
    box-shadow: 0 1px 4px rgba(0,0,0,.07);
  }}
  .metric-card .mc-label {{ font-size: 10px; text-transform: uppercase;
    letter-spacing: .8px; color: #888; margin-bottom: 4px; }}
  .metric-card .mc-value {{
    font-size: 22px; font-weight: 700; line-height: 1.1;
  }}
  .metric-card .mc-sub {{ font-size: 11px; color: #777; margin-top: 3px; }}
  .pos {{ color: var(--green); }}
  .neg {{ color: var(--red);   }}
  .neutral {{ color: #555; }}

  /* Tables */
  .section-title {{
    font-size: 13px; font-weight: 700; color: var(--dark);
    text-transform: uppercase; letter-spacing: .6px;
    border-left: 3px solid {C_BLUE}; padding-left: 8px;
    margin: 16px 0 8px;
  }}
  table.rt {{
    width: 100%; border-collapse: collapse; background: white;
    border-radius: 8px; overflow: hidden; box-shadow: 0 1px 4px rgba(0,0,0,.06);
    font-size: 12px;
  }}
  table.rt th {{
    background: var(--navy); color: white; padding: 8px 12px;
    text-align: right; font-weight: 600; font-size: 11px;
    text-transform: uppercase; letter-spacing: .5px; white-space: nowrap;
    cursor: pointer; user-select: none;
  }}
  table.rt th:hover {{ background: var(--dark); }}
  table.rt th.sort-asc::after  {{ content: " ▲"; font-size: 9px; opacity: .8; }}
  table.rt th.sort-desc::after {{ content: " ▼"; font-size: 9px; opacity: .8; }}
  table.rt th:first-child {{ text-align: left; }}
  table.rt td {{
    padding: 7px 12px; border-bottom: 1px solid #f0f0f0;
    text-align: right;
  }}
  table.rt td:first-child {{ text-align: left; font-weight: 500; }}
  table.rt tr:last-child td {{ border-bottom: none; }}
  table.rt tr.total-row td {{ font-weight: 700; background: #f8f9fb; border-top: 1px solid #ddd; }}
  table.rt tr.bench-row td {{ color: #999; font-style: italic; }}
  table.rt tr:hover td {{ background: #fafbff; }}

  /* Two-column layout */
  .two-col {{ display: grid; grid-template-columns: 1fr 1fr; gap: 16px; }}
  @media (max-width: 900px) {{ .two-col {{ grid-template-columns: 1fr; }} }}

  /* VaR bar */
  .var-bar {{
    background: white; border-radius: 8px; padding: 10px 16px;
    display: flex; align-items: center; gap: 20px; flex-wrap: wrap;
    box-shadow: 0 1px 4px rgba(0,0,0,.06); margin-top: 14px;
    font-size: 12px;
  }}
  .var-bar .vb-label {{ font-weight: 700; color: var(--dark); margin-right: 4px; }}
  .var-bar .vb-sep {{ color: #ccc; }}
  .var-util {{ display: inline-block; height: 8px; border-radius: 4px;
    background: linear-gradient(90deg, {C_BLUE}, var(--navy)); }}

  /* Market table */
  .mkt-chg.pos {{ color: var(--green); font-weight: 600; }}
  .mkt-chg.neg {{ color: var(--red);   font-weight: 600; }}
"""


# ── HTML rendering ─────────────────────────────────────────────────────────────
def _html_var_bar(fc: FundCtx) -> str:
    items = []
    for lbl, val in [("Médio", fc.var_avg_pct), ("Máx", fc.var_max_pct), ("Último", fc.var_last_pct)]:
        if val is not None:
            cls = "neg" if val > fc.var_soft else ("neutral" if val > fc.var_soft * 0.7 else "pos")
            items.append(f'<span class="vb-label">{lbl}:</span> '
                         f'<span class="{cls}">{val:.2f}%</span>')
    if fc.var_soft:
        util = (fc.var_avg_pct or 0) / fc.var_soft * 100
        items.append(f'<span class="vb-sep">|</span> '
                     f'<span class="vb-label">Limite:</span> '
                     f'Soft {fc.var_soft:.2f}% · Hard {fc.var_hard:.2f}% '
                     f'(<span class="{"neg" if util > 100 else "neutral"}">{util:.0f}% utiliz.</span>)')
    if not items:
        return ""
    return '<div class="var-bar"><span class="vb-label">VaR período</span>' + " ".join(items) + "</div>"


def _html_pa_livro_table(fc: FundCtx) -> str:
    df = fc.pa_by_livro.copy()
    if df.empty:
        return ""
    active_df = df[~df["is_bench"]]
    noise_df  = df[df["is_bench"]]
    total_mtd = float(active_df["mtd_bps"].sum())
    total_ytd = float(active_df["ytd_bps"].sum())

    rows_html = []
    for _, r in active_df.iterrows():
        pct = (r["mtd_bps"] / total_mtd * 100) if total_mtd else 0
        rows_html.append(
            f'<tr><td>{html_lib.escape(str(r["livro"]))}</td>'
            f'<td class="{_css_sign(r["mtd_bps"])}">{_fbps(r["mtd_bps"])}</td>'
            f'<td class="{_css_sign(r["ytd_bps"])}">{_fbps(r["ytd_bps"])}</td>'
            f'<td class="neutral">{pct:+.0f}%</td></tr>'
        )
    rows_html.append(
        f'<tr class="total-row"><td>Total ativo</td>'
        f'<td class="{_css_sign(total_mtd)}">{_fbps(total_mtd)}</td>'
        f'<td class="{_css_sign(total_ytd)}">{_fbps(total_ytd)}</td>'
        f'<td>100%</td></tr>'
    )
    for _, r in noise_df.iterrows():
        rows_html.append(
            f'<tr class="bench-row"><td>{html_lib.escape(str(r["livro"]))}</td>'
            f'<td class="neutral">{_fbps(r["mtd_bps"])}</td>'
            f'<td class="neutral">{_fbps(r["ytd_bps"])}</td>'
            f'<td>—</td></tr>'
        )
    return (
        '<table class="rt">'
        '<thead><tr><th>Categoria</th><th>MTD α</th><th>YTD α</th><th>% MTD</th></tr></thead>'
        '<tbody>' + "".join(rows_html) + "</tbody></table>"
    )


def _html_pa_products_tables(fc: FundCtx) -> str:
    df = fc.pa_products
    if df.empty:
        return ""

    def _tbl(rows: pd.DataFrame, col: str) -> str:
        if rows.empty:
            return ""
        trs = ""
        for _, r in rows.iterrows():
            trs += (
                f'<tr><td>{html_lib.escape(str(r["product"]))}</td>'
                f'<td style="color:#888;font-size:10px">{html_lib.escape(str(r["livro"]))}</td>'
                f'<td class="{_css_sign(r["mtd_bps"])}">{_fbps(r["mtd_bps"])}</td>'
                f'<td class="{_css_sign(r["ytd_bps"])}">{_fbps(r["ytd_bps"])}</td></tr>'
            )
        return (
            '<table class="rt">'
            '<thead><tr><th>Instrumento</th><th>Livro</th><th>MTD α</th><th>YTD α</th></tr></thead>'
            f'<tbody>{trs}</tbody></table>'
        )

    mtd_top = df.nlargest(5, "mtd_bps")
    mtd_bot = df.nsmallest(5, "mtd_bps")
    ytd_top = df.nlargest(5, "ytd_bps")
    ytd_bot = df.nsmallest(5, "ytd_bps")

    return (
        '<div class="section-title">Top Contribuidores / Detratores — MTD</div>'
        '<div class="two-col">'
        + _tbl(mtd_top, "mtd_bps") + _tbl(mtd_bot, "mtd_bps")
        + '</div>'
        '<div class="section-title" style="margin-top:14px">Top Contribuidores / Detratores — YTD</div>'
        '<div class="two-col">'
        + _tbl(ytd_top, "ytd_bps") + _tbl(ytd_bot, "ytd_bps")
        + '</div>'
    )


def _html_fund_tab(fc: FundCtx, period_label: str, peers: list[dict] | None = None) -> str:
    alpha_sub = (f"vs {fc.bench_name} MTD: {_fbps(fc.mtd_alpha_bps - fc.bench_mtd_bps)}"
                 if fc.bench_mtd_bps else "")

    cards = [
        ("Alpha MTD (líq. bench)", _fbps(fc.mtd_alpha_bps), _css_sign(fc.mtd_alpha_bps), alpha_sub),
        ("Alpha YTD (líq. bench)", _fbps(fc.ytd_alpha_bps), _css_sign(fc.ytd_alpha_bps),
         f"{fc.bench_name} YTD: {_fbps(fc.bench_ytd_bps)}"),
    ]
    if fc.fund_mtd_pct is not None:
        cards.append((
            "Retorno Bruto MTD", _fpct(fc.fund_mtd_pct), _css_sign(fc.fund_mtd_pct),
            f"YTD: {_fpct(fc.fund_ytd_pct)}" if fc.fund_ytd_pct is not None else ""
        ))

    cards_html = "".join(
        f'<div class="metric-card">'
        f'<div class="mc-label">{lbl}</div>'
        f'<div class="mc-value {cls}">{val}</div>'
        f'<div class="mc-sub">{sub}</div>'
        f'</div>'
        for lbl, val, cls, sub in cards
    )

    return (
        f'<div class="metric-row">{cards_html}</div>'
        f'<div class="section-title">Performance Attribution — {period_label}</div>'
        + _html_pa_livro_table(fc)
        + _html_pa_products_tables(fc)
        + _html_var_bar(fc)
        + (_html_peers_section(peers, fund_label=fc.label) if peers else "")
    )


def _html_idka_tab(fc_3y: FundCtx, fc_10y: FundCtx, period_label: str) -> str:
    def _half(fc: FundCtx) -> str:
        return (
            f'<h3 style="color:var(--dark);font-size:14px;margin-bottom:10px">{fc.label}</h3>'
            + _html_fund_tab(fc, period_label)
        )
    return (
        '<div class="two-col">'
        + _half(fc_3y) + _half(fc_10y)
        + "</div>"
    )


def _html_market_tab(market_rows: list[dict], period_label: str) -> str:
    if not market_rows:
        return '<p style="color:#999;padding:20px">Dados de mercado não disponíveis.</p>'
    trs = ""
    for r in market_rows:
        chg_mtd_cls = "pos" if r["chg_mtd"] >= 0 else "neg"
        chg_ytd_cls = "pos" if r["chg_ytd"] >= 0 else "neg"
        unit = r["unit"]
        trs += (
            f'<tr><td>{html_lib.escape(r["label"])}</td>'
            f'<td class="neutral">{r["val_start"]:.2f}</td>'
            f'<td class="neutral">{r["val_end"]:.2f}</td>'
            f'<td class="mkt-chg {chg_mtd_cls}">{r["chg_mtd"]:+.2f} {unit}</td>'
            f'<td class="mkt-chg {chg_ytd_cls}">{r["chg_ytd"]:+.2f} {unit}</td>'
            f'</tr>'
        )
    return (
        f'<div class="section-title">Mercado — {period_label}</div>'
        '<table class="rt" style="max-width:700px">'
        f'<thead><tr><th>Indicador</th><th>Início período</th>'
        f'<th>Final período</th><th>Var MTD</th><th>Var YTD</th></tr></thead>'
        f'<tbody>{trs}</tbody></table>'
    )


def build_html(
    month_label: str, period_str: str, is_final: bool,
    fund_contexts: dict[str, FundCtx],
    market_rows: list[dict],
    peers_data: dict | None = None,
) -> str:
    status = "Final" if is_final else "Parcial"
    tab_btns = '<button class="tab-btn active" onclick="showTab(\'market\')">Mercado</button>'
    tab_divs = f'<div id="tab-market" class="tab-content active">{_html_market_tab(market_rows, month_label)}</div>'

    for tab_id in _TAB_ORDER:
        if tab_id == "IDKA":
            fc_3y  = fund_contexts.get("IDKA_3Y")
            fc_10y = fund_contexts.get("IDKA_10Y")
            if fc_3y is None and fc_10y is None:
                continue
            content = _html_idka_tab(
                fc_3y  or FundCtx("IDKA_3Y",  "IDKA 3Y",  "IDKA 3Y",  pd.DataFrame(), pd.DataFrame()),
                fc_10y or FundCtx("IDKA_10Y", "IDKA 10Y", "IDKA 10Y", pd.DataFrame(), pd.DataFrame()),
                month_label,
            )
        else:
            fc = fund_contexts.get(tab_id)
            if fc is None:
                continue
            peers = _get_peer_list(tab_id, peers_data or {}) if peers_data else []
            content = _html_fund_tab(fc, month_label, peers=peers)

        lbl = _TAB_LABELS[tab_id]
        tab_btns += f'<button class="tab-btn" onclick="showTab(\'{tab_id}\')">{lbl}</button>'
        tab_divs += f'<div id="tab-{tab_id}" class="tab-content">{content}</div>'

    js = """
    function showTab(id) {
      document.querySelectorAll('.tab-content').forEach(el => el.classList.remove('active'));
      document.querySelectorAll('.tab-btn').forEach(el => el.classList.remove('active'));
      document.getElementById('tab-' + id).classList.add('active');
      event.currentTarget.classList.add('active');
    }

    // Universal table sort — click any <th> in table.rt
    (function() {
      function parseVal(s) {
        // strip %, +, bps, spaces, commas; parse as number if possible
        var n = parseFloat(s.replace(/[+%,\s]/g,'').replace('bps','').replace('—',''));
        return isNaN(n) ? s.toLowerCase() : n;
      }
      document.addEventListener('click', function(e) {
        var th = e.target.closest('th');
        if (!th) return;
        var tbl = th.closest('table.rt');
        if (!tbl) return;
        var ths = Array.from(th.parentElement.children);
        var col = ths.indexOf(th);
        var asc = !th.classList.contains('sort-asc');
        ths.forEach(function(h){ h.classList.remove('sort-asc','sort-desc'); });
        th.classList.add(asc ? 'sort-asc' : 'sort-desc');
        var tbody = tbl.querySelector('tbody');
        var rows = Array.from(tbody.querySelectorAll('tr')).filter(function(r){
          return !r.classList.contains('total-row');
        });
        var pinned = Array.from(tbody.querySelectorAll('tr.total-row'));
        rows.sort(function(a, b) {
          var av = parseVal(a.cells[col] ? a.cells[col].textContent : '');
          var bv = parseVal(b.cells[col] ? b.cells[col].textContent : '');
          if (av < bv) return asc ? -1 : 1;
          if (av > bv) return asc ? 1 : -1;
          return 0;
        });
        rows.forEach(function(r){ tbody.appendChild(r); });
        pinned.forEach(function(r){ tbody.appendChild(r); });
      });
    })();
    """

    return f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Monthly Review — {month_label}</title>
  <style>{_CSS}</style>
</head>
<body>
<div class="header">
  <div>
    <div class="header-title">Galapagos Capital · Monthly Review</div>
    <div class="header-sub">{month_label} &nbsp;·&nbsp; {period_str}</div>
  </div>
  <div class="header-badge">{status}</div>
</div>
<div class="tab-bar">{tab_btns}</div>
<div style="max-width:1300px;margin:0 auto">
  {tab_divs}
</div>
<script>{js}</script>
</body>
</html>"""


# ── PPTX helpers ──────────────────────────────────────────────────────────────
def _pptx_add_textbox(slide, left, top, width, height, text,
                       font_size=12, bold=False, italic=False,
                       color="000000", align="left", wrap=True):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    return txBox


def _pptx_add_rect(slide, left, top, width, height,
                    fill_hex: str | None = None, line_hex: str | None = None,
                    line_width_pt: float = 0.75):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    else:
        shape.fill.background()
    if line_hex:
        shape.line.color.rgb = RGBColor.from_string(line_hex)
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def _pptx_bar_chart_image(labels: list[str], values: list[float],
                            title: str = "", width_in: float = 5.5,
                            height_in: float = 3.0) -> BytesIO:
    """Horizontal bar chart: green=positive, red=negative. Returns BytesIO PNG."""
    n = len(labels)
    colors = [C_GREEN if v >= 0 else C_RED for v in values]

    fig, ax = plt.subplots(figsize=(width_in, max(height_in, n * 0.38 + 0.5)))
    fig.patch.set_facecolor("white")
    bars = ax.barh(range(n), values, color=colors, height=0.55, zorder=3)
    ax.set_yticks(range(n))
    ax.set_yticklabels(labels, fontsize=9)
    ax.invert_yaxis()
    ax.axvline(0, color="#aaa", linewidth=0.8)
    ax.set_xlabel("bps (alpha líq. bench)", fontsize=8)
    if title:
        ax.set_title(title, fontsize=10, fontweight="bold", color=C_DARK)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="y", length=0)
    ax.set_facecolor("white")

    for bar, val in zip(bars, values):
        x_pos = val + (max(abs(v) for v in values) * 0.02 if values else 0)
        x_pos = val + 1 if val >= 0 else val - 1
        ha = "left" if val >= 0 else "right"
        ax.text(x_pos, bar.get_y() + bar.get_height() / 2,
                f"{val:+.1f}", va="center", ha=ha, fontsize=8)

    plt.tight_layout(pad=0.5)
    buf = BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)
    return buf


def _pptx_metric_cards(slide, fc: FundCtx, top: float):
    """Add 3 metric boxes to slide at given top position."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    cards = [
        ("Alpha MTD\n(líq. bench)", _fbps(fc.mtd_alpha_bps),
         "1E8C45" if fc.mtd_alpha_bps >= 0 else "C0392B"),
        ("Alpha YTD\n(líq. bench)", _fbps(fc.ytd_alpha_bps),
         "1E8C45" if fc.ytd_alpha_bps >= 0 else "C0392B"),
        (f"Retorno MTD\nvs {fc.bench_name}",
         (_fpct(fc.fund_mtd_pct) if fc.fund_mtd_pct is not None else "—") +
         f"\n{fc.bench_name} {_fbps(fc.bench_mtd_bps)}",
         "183B80"),
    ]
    card_w, card_h = 3.8, 0.85
    for i, (lbl, val, color_hex) in enumerate(cards):
        x = 0.2 + i * (card_w + 0.1)
        _pptx_add_rect(slide, x, top, card_w, card_h, fill_hex="F5F7FA",
                        line_hex="DEE2E9", line_width_pt=0.5)
        _pptx_add_textbox(slide, x + 0.12, top + 0.06, card_w - 0.2, 0.28,
                           lbl, font_size=8, color="888888")
        _pptx_add_textbox(slide, x + 0.12, top + 0.33, card_w - 0.2, 0.42,
                           val, font_size=13, bold=True, color=color_hex)


def _pptx_comment_box(slide, left: float, top: float, width: float, height: float, fund_label: str):
    """Add empty comment placeholder for PM to fill in."""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    _pptx_add_rect(slide, left, top, width, height,
                    fill_hex="F0F2F5", line_hex="183B80", line_width_pt=0.75)
    _pptx_add_textbox(slide, left + 0.12, top + 0.10, width - 0.2, 0.28,
                       "COMENTÁRIOS DO GESTOR", font_size=8, bold=True,
                       color="183B80")
    _pptx_add_textbox(slide, left + 0.12, top + 0.42, width - 0.2, height - 0.55,
                       "(adicione aqui a análise qualitativa do período)",
                       font_size=9, italic=True, color="AAAAAA")


def _pptx_var_line(slide, fc: FundCtx, top: float):
    parts = []
    for lbl, v in [("Médio", fc.var_avg_pct), ("Máx", fc.var_max_pct), ("Último", fc.var_last_pct)]:
        if v is not None:
            parts.append(f"VaR {lbl}: {v:.2f}%")
    if fc.var_soft:
        util = (fc.var_avg_pct or 0) / fc.var_soft * 100
        parts.append(f"Limite Soft: {fc.var_soft:.2f}%  Hard: {fc.var_hard:.2f}%  ({util:.0f}% utiliz. médio)")
    if parts:
        _pptx_add_rect(slide, 0.2, top, 12.93, 0.45, fill_hex="EEF0F4", line_hex="DEE2E9")
        _pptx_add_textbox(slide, 0.35, top + 0.08, 12.6, 0.30,
                           "  ·  ".join(parts), font_size=9, color="133979")


def _pptx_fund_slide(prs, fc: FundCtx, month_label: str):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # Header bar
    _pptx_add_rect(slide, 0, 0, 13.33, 0.65, fill_hex="183B80")
    _pptx_add_textbox(slide, 0.25, 0.08, 9.5, 0.50,
                       fc.label.upper(), font_size=20, bold=True, color="FFFFFF")
    _pptx_add_textbox(slide, 9.8, 0.12, 3.3, 0.40,
                       month_label, font_size=12, color="AACFFF", align="right")

    # Metric cards
    _pptx_metric_cards(slide, fc, top=0.72)

    # PA bar chart (left 55%)
    active = fc.pa_by_livro[~fc.pa_by_livro["is_bench"]].copy()
    active = active.sort_values("mtd_bps")
    if not active.empty:
        chart_buf = _pptx_bar_chart_image(
            labels=list(active["livro"]),
            values=list(active["mtd_bps"]),
            title=f"PA por Estratégia — MTD α (líq. bench)",
            width_in=5.8, height_in=3.2,
        )
        slide.shapes.add_picture(chart_buf, Inches(0.2), Inches(1.68),
                                  Inches(7.0), Inches(3.3))

    # Comment box (right, compact — PM fills in)
    _pptx_comment_box(slide, 7.4, 1.68, 5.7, 1.6, fc.label)

    # Top contributors mini-table as text
    top5 = fc.pa_products.nlargest(5, "mtd_bps")
    bot5 = fc.pa_products.nsmallest(5, "mtd_bps")

    def _contrib_text(df: pd.DataFrame) -> str:
        lines = []
        for _, r in df.iterrows():
            lines.append(f"{str(r['product'])[:28]:28s}  {r['mtd_bps']:+.1f} bps")
        return "\n".join(lines)

    if not top5.empty:
        _pptx_add_textbox(slide, 0.2, 5.1, 6.2, 0.22,
                           "TOP 5 CONTRIBUIDORES (MTD)", font_size=8, bold=True, color="1E8C45")
        _pptx_add_textbox(slide, 0.2, 5.35, 6.2, 1.2,
                           _contrib_text(top5), font_size=8, color="333333")
    if not bot5.empty:
        _pptx_add_textbox(slide, 6.5, 5.1, 6.6, 0.22,
                           "TOP 5 DETRATORES (MTD)", font_size=8, bold=True, color="C0392B")
        _pptx_add_textbox(slide, 6.5, 5.35, 6.6, 1.2,
                           _contrib_text(bot5), font_size=8, color="333333")

    # VaR line
    _pptx_var_line(slide, fc, top=6.65)


def _pptx_idka_slide(prs, fc_3y: FundCtx, fc_10y: FundCtx, month_label: str):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    _pptx_add_rect(slide, 0, 0, 13.33, 0.65, fill_hex="183B80")
    _pptx_add_textbox(slide, 0.25, 0.08, 9.5, 0.50,
                       "IDKA — IPCA 3Y & 10Y", font_size=20, bold=True, color="FFFFFF")
    _pptx_add_textbox(slide, 9.8, 0.12, 3.3, 0.40,
                       month_label, font_size=12, color="AACFFF", align="right")

    # Vertical divider
    _pptx_add_rect(slide, 6.55, 0.68, 0.04, 6.7, fill_hex="DEE2E9")

    for i, fc in enumerate([fc_3y, fc_10y]):
        x_off = 0.0 if i == 0 else 6.65
        col_w = 6.4

        # Fund name
        _pptx_add_textbox(slide, x_off + 0.2, 0.72, col_w, 0.32,
                           fc.label, font_size=13, bold=True, color="133979")

        # Metric cards (compact 2-wide)
        for j, (lbl, val, col_hex) in enumerate([
            ("Alpha MTD", _fbps(fc.mtd_alpha_bps),
             "1E8C45" if fc.mtd_alpha_bps >= 0 else "C0392B"),
            ("Alpha YTD", _fbps(fc.ytd_alpha_bps),
             "1E8C45" if fc.ytd_alpha_bps >= 0 else "C0392B"),
        ]):
            cx = x_off + 0.2 + j * 3.1
            _pptx_add_rect(slide, cx, 1.10, 2.85, 0.70, fill_hex="F5F7FA",
                            line_hex="DEE2E9", line_width_pt=0.5)
            _pptx_add_textbox(slide, cx + 0.10, 1.15, 2.6, 0.22,
                               lbl, font_size=8, color="888888")
            _pptx_add_textbox(slide, cx + 0.10, 1.38, 2.6, 0.35,
                               val, font_size=12, bold=True, color=col_hex)

        # PA bar chart
        active = fc.pa_by_livro[~fc.pa_by_livro["is_bench"]].sort_values("mtd_bps")
        if not active.empty:
            chart_buf = _pptx_bar_chart_image(
                labels=list(active["livro"]), values=list(active["mtd_bps"]),
                title="", width_in=5.2, height_in=2.0,
            )
            slide.shapes.add_picture(chart_buf, Inches(x_off + 0.15), Inches(1.92),
                                      Inches(6.2), Inches(2.1))

        # VaR stats
        vs_parts = []
        for lbl, v in [("Médio", fc.var_avg_pct), ("Máx", fc.var_max_pct), ("Último", fc.var_last_pct)]:
            if v is not None:
                vs_parts.append(f"BVaR {lbl}: {v:.2f}%")
        if fc.var_soft and vs_parts:
            vs_parts.append(f"Limite: {fc.var_soft:.2f}%")
        if vs_parts:
            _pptx_add_textbox(slide, x_off + 0.2, 4.08, col_w, 0.28,
                               "  ·  ".join(vs_parts), font_size=8, color="555555")

        # Comment box (compact)
        _pptx_comment_box(slide, x_off + 0.2, 4.42, col_w - 0.1, 1.5, fc.label)


def _pptx_cover_slide(prs, month_label: str, period_str: str, is_final: bool):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_rect(slide, 0, 0, 13.33, 7.50, fill_hex="183B80")
    _pptx_add_rect(slide, 0, 5.2, 13.33, 2.3, fill_hex="0C2048")
    _pptx_add_rect(slide, 0.35, 2.85, 4.5, 0.06, fill_hex="2AADF5")
    _pptx_add_textbox(slide, 0.4, 1.5, 12.0, 1.0,
                       "GALAPAGOS CAPITAL", font_size=14, bold=True, color="2AADF5")
    _pptx_add_textbox(slide, 0.4, 2.2, 12.0, 1.1,
                       "Monthly Review", font_size=36, bold=True, color="FFFFFF")
    _pptx_add_textbox(slide, 0.4, 3.15, 8.0, 0.6,
                       month_label, font_size=24, bold=False, color="AACFFF")
    _pptx_add_textbox(slide, 0.4, 3.85, 8.0, 0.4,
                       period_str + (" · Final" if is_final else " · Parcial"),
                       font_size=12, color="7FAADD")
    _pptx_add_textbox(slide, 0.4, 6.0, 12.0, 0.5,
                       "Galapagos Capital · Asset Management · Confidencial",
                       font_size=10, color="6688AA")


def _pptx_market_slide(prs, market_rows: list[dict], month_label: str):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_rect(slide, 0, 0, 13.33, 0.65, fill_hex="183B80")
    _pptx_add_textbox(slide, 0.25, 0.08, 9.5, 0.50,
                       "VISÃO DE MERCADO", font_size=20, bold=True, color="FFFFFF")
    _pptx_add_textbox(slide, 9.8, 0.12, 3.3, 0.40,
                       month_label, font_size=12, color="AACFFF", align="right")

    if not market_rows:
        _pptx_add_textbox(slide, 0.4, 1.0, 12.0, 0.5,
                           "Dados de mercado não disponíveis.", font_size=12, color="999999")
        return

    headers = ["Indicador", "Início período", "Final período", "Var MTD", "Var YTD"]
    col_ws  = [2.5, 1.8, 1.8, 1.8, 1.8]
    x_starts = [0.2]
    for w in col_ws[:-1]:
        x_starts.append(x_starts[-1] + w + 0.05)

    # Header row
    y = 0.82
    _pptx_add_rect(slide, 0.2, y, sum(col_ws) + 0.2, 0.35, fill_hex="183B80")
    for h, x, w in zip(headers, x_starts, col_ws):
        _pptx_add_textbox(slide, x + 0.08, y + 0.06, w, 0.24,
                           h, font_size=9, bold=True, color="FFFFFF")

    for i, r in enumerate(market_rows):
        y = 1.22 + i * 0.42
        bg = "F5F7FA" if i % 2 == 0 else "FFFFFF"
        _pptx_add_rect(slide, 0.2, y, sum(col_ws) + 0.2, 0.38, fill_hex=bg,
                        line_hex="DEE2E9", line_width_pt=0.4)
        vals = [
            r["label"],
            f"{r['val_start']:.2f}",
            f"{r['val_end']:.2f}",
            f"{r['chg_mtd']:+.2f} {r['unit']}",
            f"{r['chg_ytd']:+.2f} {r['unit']}",
        ]
        for j, (val, x, w) in enumerate(zip(vals, x_starts, col_ws)):
            color = "222222"
            if j == 3:
                color = "1E8C45" if r["chg_mtd"] >= 0 else "C0392B"
            elif j == 4:
                color = "1E8C45" if r["chg_ytd"] >= 0 else "C0392B"
            _pptx_add_textbox(slide, x + 0.08, y + 0.09, w, 0.24,
                               val, font_size=10, color=color,
                               bold=(j == 3 or j == 4))

    # Comment box for macro context
    comment_y = 1.22 + len(market_rows) * 0.42 + 0.2
    if comment_y < 6.5:
        _pptx_add_rect(slide, 0.2, comment_y, 12.9, 7.38 - comment_y - 0.1,
                        fill_hex="F0F2F5", line_hex="183B80", line_width_pt=0.75)
        _pptx_add_textbox(slide, 0.35, comment_y + 0.12, 12.5, 0.28,
                           "CONTEXTO MACRO — COMENTÁRIO DO PERÍODO", font_size=8,
                           bold=True, color="183B80")
        _pptx_add_textbox(slide, 0.35, comment_y + 0.45, 12.5, 7.38 - comment_y - 0.6,
                           "(adicione aqui a narrativa macro do período)",
                           font_size=9, italic=True, color="AAAAAA")


def _pptx_peers_slide(prs, fc: FundCtx, peers: list[dict], month_label: str):
    """Peers comparison slide: 4 distribution bars (MTD+YTD) + 2 scatter charts."""
    from pptx.util import Inches

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header — slightly different shade to distinguish from main slide
    _pptx_add_rect(slide, 0, 0, 13.33, 0.65, fill_hex="133979")
    _pptx_add_textbox(slide, 0.25, 0.08, 9.5, 0.50,
                       f"{fc.label.upper()} — PEERS", font_size=20, bold=True, color="FFFFFF")
    _pptx_add_textbox(slide, 9.8, 0.12, 3.3, 0.40,
                       month_label, font_size=12, color="AACFFF", align="right")

    buf = _peers_chart(peers, fc.label)
    if buf.getbuffer().nbytes > 100:
        slide.shapes.add_picture(buf, Inches(0.15), Inches(0.72),
                                  Inches(12.95), Inches(6.65))


def _pptx_closing_slide(prs, month_label: str):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _pptx_add_rect(slide, 0, 0, 13.33, 7.50, fill_hex="0C2048")
    _pptx_add_rect(slide, 0, 6.8, 13.33, 0.70, fill_hex="183B80")
    _pptx_add_textbox(slide, 0.4, 2.5, 12.0, 0.8,
                       "galapagoscapital.com", font_size=20, bold=True, color="2AADF5")
    _pptx_add_textbox(slide, 0.4, 3.5, 12.0, 1.5,
                       "Este documento contém informações confidenciais destinadas exclusivamente "
                       "ao uso interno da Galapagos Capital. Não deve ser distribuído sem "
                       "autorização prévia.",
                       font_size=9, color="7799BB", wrap=True)
    _pptx_add_textbox(slide, 0.4, 6.85, 12.0, 0.5,
                       f"Monthly Review · {month_label}", font_size=10, color="AACFFF")


def build_pptx(
    month_label: str, period_str: str, is_final: bool,
    fund_contexts: dict[str, FundCtx],
    market_rows: list[dict],
    peers_data: dict | None = None,
) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.50)

    _pptx_cover_slide(prs, month_label, period_str, is_final)
    _pptx_market_slide(prs, market_rows, month_label)

    for short in FUND_ORDER:
        if short in ("IDKA_3Y", "IDKA_10Y"):
            continue  # handled as combined slide below
        fc = fund_contexts.get(short)
        if fc is None:
            continue
        _pptx_fund_slide(prs, fc, month_label)
        # Peers slide immediately after (only for funds with peer groups)
        peers = _get_peer_list(short, peers_data or {})
        if peers:
            _pptx_peers_slide(prs, fc, peers, month_label)

    # Combined IDKA slide (no peers for IDKAs)
    fc_3y  = fund_contexts.get("IDKA_3Y")
    fc_10y = fund_contexts.get("IDKA_10Y")
    if fc_3y or fc_10y:
        _pptx_idka_slide(
            prs,
            fc_3y  or FundCtx("IDKA_3Y",  "IDKA 3Y",  "IDKA 3Y",  pd.DataFrame(), pd.DataFrame()),
            fc_10y or FundCtx("IDKA_10Y", "IDKA 10Y", "IDKA 10Y", pd.DataFrame(), pd.DataFrame()),
            month_label,
        )

    _pptx_closing_slide(prs, month_label)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Main ──────────────────────────────────────────────────────────────────────
def main() -> None:
    parser = argparse.ArgumentParser(description="Monthly Review Report — GLPG")
    parser.add_argument("--month", metavar="YYYY-MM", default=None,
                        help="Month to generate (default: prompts)")
    args = parser.parse_args()

    if args.month:
        ym = args.month.strip()
    else:
        default_ym = (pd.Timestamp("today") - pd.offsets.MonthBegin(1)).strftime("%Y-%m")
        ym_input   = input(f"Mês (YYYY-MM) [{default_ym}]: ").strip()
        ym         = ym_input if ym_input else default_ym

    # Validate format
    import re
    if not re.fullmatch(r"\d{4}-\d{2}", ym):
        sys.exit(f"Formato inválido: '{ym}'. Use YYYY-MM.")

    print(f"\n▶ Resolvendo datas para {ym}...")
    bdays, start, end, year_start, is_final, month_label = _resolve_dates(ym)
    period_str = f"{pd.Timestamp(start).strftime('%d/%m')} – {pd.Timestamp(end).strftime('%d/%m/%Y')}"
    status = "Final" if is_final else f"Parcial (até {pd.Timestamp(end).strftime('%d/%m')})"
    print(f"  Período: {period_str}  ·  {status}")
    print(f"  Dias úteis: {len(bdays)}")

    # ── Concurrent data fetch ─────────────────────────────────────────────────
    print("\n▶ Buscando dados...")
    results: dict = {}
    tasks = {
        "pa":      lambda: _fetch_pa(start, end, year_start),
        "nav":     lambda: _fetch_nav_series(year_start, end),
        "cdi":     lambda: _fetch_cdi(year_start, end, start),
        "ibov":    lambda: _fetch_ibov(year_start, end, start),
        "idka":    lambda: _fetch_idka_returns(year_start, end, start),
        "market":  lambda: _fetch_market_moves(start, end, year_start),
    }

    with ThreadPoolExecutor(max_workers=6) as executor:
        futures = {executor.submit(fn): key for key, fn in tasks.items()}
        for future in as_completed(futures):
            key = futures[future]
            try:
                results[key] = future.result()
                print(f"  ✓ {key}")
            except Exception as exc:
                print(f"  ✗ {key}: {exc}")
                results[key] = None

    _r = results.get
    pa_df       = _r("pa")  if _r("pa")  is not None else pd.DataFrame()
    nav_df      = _r("nav") if _r("nav") is not None else pd.DataFrame()
    market_rows = _r("market") if _r("market") is not None else []

    benchmarks = {
        "cdi":  _r("cdi")  if _r("cdi")  is not None else {"mtd_bps": 0.0, "ytd_bps": 0.0},
        "ibov": _r("ibov") if _r("ibov") is not None else {"mtd_bps": 0.0, "ytd_bps": 0.0},
        "idka": _r("idka") if _r("idka") is not None else {},
    }

    print("\n▶ Calculando VaR do período...")
    var_stats: dict[str, dict] = {}
    if not nav_df.empty:
        try:
            var_stats = _fetch_var_period(start, end, nav_df)
            print(f"  ✓ var_stats para {list(var_stats.keys())}")
        except Exception as exc:
            print(f"  ✗ var_stats: {exc}")

    # ── Build per-fund contexts ───────────────────────────────────────────────
    print("\n▶ Montando contextos por fundo...")
    fund_contexts: dict[str, FundCtx] = {}
    for short in FUND_ORDER:
        try:
            fc = _build_fund_ctx(short, pa_df, nav_df, benchmarks, var_stats,
                                  start, end, year_start)
            fund_contexts[short] = fc
            print(f"  ✓ {short}  MTD α {_fbps(fc.mtd_alpha_bps)}  YTD α {_fbps(fc.ytd_alpha_bps)}")
        except Exception as exc:
            print(f"  ✗ {short}: {exc}")

    # ── Load peers data (self-contained: tries local copy first, then shared drive) ──
    print("\n▶ Carregando dados de peers...")
    try:
        peers_data = _load_peers_data()
        n_groups = len(peers_data.get("groups", {}))
        print(f"  ✓ {n_groups} grupos de peers  (data: {peers_data.get('val_date', '?')})")
    except Exception as exc:
        print(f"  ✗ peers: {exc}")
        peers_data = {}

    # ── Render HTML ───────────────────────────────────────────────────────────
    print("\n▶ Gerando HTML...")
    html_str = build_html(month_label, period_str, is_final, fund_contexts,
                          market_rows, peers_data=peers_data)
    html_path = OUT_DIR / f"{ym}_monthly_review.html"
    html_path.write_text(html_str, encoding="utf-8")
    print(f"  → {html_path}")

    # ── Render PPTX ───────────────────────────────────────────────────────────
    print("\n▶ Gerando PPTX...")
    pptx_bytes = build_pptx(month_label, period_str, is_final, fund_contexts,
                             market_rows, peers_data=peers_data)
    pptx_path = OUT_DIR / f"{ym}_monthly_review.pptx"
    pptx_path.write_bytes(pptx_bytes)
    print(f"  → {pptx_path}")

    # ── Mirror to shared distribution location ────────────────────────────────
    # Two file pairs written to share:
    #   {ym}_monthly_review.{html,pptx}   — dated archive (acumula histórico)
    #   ultimo_monthly_review.{html,pptx} — "latest" pointer (nome fixo,
    #                                       sempre sobrescreve)
    mirror_dir = Path(r"F:\Bloomberg\Risk_Manager\Data\Monthlyreview")
    try:
        mirror_dir.mkdir(parents=True, exist_ok=True)
        (mirror_dir / html_path.name).write_text(html_str, encoding="utf-8")
        (mirror_dir / pptx_path.name).write_bytes(pptx_bytes)
        print(f"  → mirror: {mirror_dir}")

        latest_html = mirror_dir / "ultimo_monthly_review.html"
        latest_pptx = mirror_dir / "ultimo_monthly_review.pptx"
        latest_html.write_text(html_str, encoding="utf-8")
        latest_pptx.write_bytes(pptx_bytes)
        print(f"  → ultimo: {latest_html.name}, {latest_pptx.name}")
    except Exception as exc:
        print(f"  ⚠ mirror falhou: {exc}")

    print(f"\n✅ Monthly Review {month_label} concluído ({status})\n")


if __name__ == "__main__":
    main()
